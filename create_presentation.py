from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(26, 26, 46)
ACCENT_BLUE = RGBColor(0, 212, 255)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(180, 180, 180)
GREEN = RGBColor(34, 197, 94)
YELLOW = RGBColor(234, 179, 8)
CARD_BG = RGBColor(40, 40, 70)
CARD_BG2 = RGBColor(30, 30, 55)

def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    sp = bg._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = LIGHT_GRAY
        p2.alignment = PP_ALIGN.CENTER

def add_section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3), Inches(2.333), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.333), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    cb = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
    tf = cb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = f"• {item}" if item else ""
        para.font.size = Pt(22)
        para.font.color.rgb = WHITE
        para.space_after = Pt(10)

def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5)).table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(16)
        para.font.color.rgb = ACCENT_BLUE

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG2
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14)
            para.font.color.rgb = GREEN if j == 1 else WHITE

def add_provider_slide(name, tagline, strengths, tradeoff, models):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = name
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    tg = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.5))
    p = tg.text_frame.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(22)
    p.font.italic = True
    p.font.color.rgb = LIGHT_GRAY

    sb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(7), Inches(3))
    tf = sb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Strengths:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    for s in strengths:
        para = tf.add_paragraph()
        para.text = f"✓ {s}"
        para.font.size = Pt(18)
        para.font.color.rgb = WHITE
        para.space_after = Pt(6)

    tw = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5), Inches(7), Inches(0.7))
    tw.fill.solid()
    tw.fill.fore_color.rgb = RGBColor(60, 50, 30)
    tw.line.fill.background()

    twb = slide.shapes.add_textbox(Inches(0.7), Inches(5.15), Inches(6.6), Inches(0.4))
    p = twb.text_frame.paragraphs[0]
    p.text = f"Trade-off: {tradeoff}"
    p.font.size = Pt(16)
    p.font.color.rgb = YELLOW

    mb = slide.shapes.add_textbox(Inches(8), Inches(1.8), Inches(4.8), Inches(4))
    tf = mb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Models:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    for m in models:
        para = tf.add_paragraph()
        para.text = m
        para.font.size = Pt(17)
        para.font.color.rgb = LIGHT_GRAY
        para.space_after = Pt(6)

def add_usecase_slide(title, recs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    y = 1.3
    for r in recs:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.333), Inches(1.05))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.fill.background()

        nb = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.12), Inches(3.5), Inches(0.4))
        p = nb.text_frame.paragraphs[0]
        p.text = r["need"]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = LIGHT_GRAY

        ab = slide.shapes.add_textbox(Inches(4.2), Inches(y + 0.12), Inches(0.5), Inches(0.4))
        p = ab.text_frame.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(18)
        p.font.color.rgb = ACCENT_BLUE

        mmb = slide.shapes.add_textbox(Inches(4.7), Inches(y + 0.12), Inches(4), Inches(0.4))
        p = mmb.text_frame.paragraphs[0]
        p.text = r["model"]
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = GREEN

        wb = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.52), Inches(11.8), Inches(0.4))
        p = wb.text_frame.paragraphs[0]
        p.text = r["why"]
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY

        y += 1.2

# ========== SLIDES ==========

add_title_slide("LLM Selection Guide", "Finding the right AI model for your project")

add_content_slide("What We'll Cover", [
    "Quick decision guide - find your LLM in seconds",
    "Provider overview - Claude, OpenAI, Gemini, DeepSeek, Open Source",
    "Use case recommendations - which model for which project",
    "Cost comparison - budget vs mid-range vs premium",
    "Resources for deeper research"
])

add_section_slide("Quick Decision Guide")

add_table_slide("What Matters Most to You?",
    ["Need", "Best Choice", "Runner Up", "Budget Option"],
    [
        ["Lowest Cost", "DeepSeek V3.2", "Gemini Flash", "Claude Haiku"],
        ["Highest Quality", "Claude Opus 4.5", "GPT-5", "Gemini 3 Pro"],
        ["Fastest Speed", "Claude Haiku", "Gemini Flash", "GPT-4o-mini"],
        ["Data Privacy", "Llama 4 (self-host)", "Qwen 3", "DeepSeek"],
        ["Best Coding", "Claude Opus 4.5", "Gemini Flash", "DeepSeek Coder"],
        ["Complex Reasoning", "OpenAI o3", "DeepSeek R1", "Claude Opus"],
    ]
)

add_section_slide("Provider Overview")

add_provider_slide("Claude (Anthropic)",
    "Best for coding & following instructions",
    ["Excellent at complex, detailed instructions",
     "Industry-leading code generation",
     "Natural, human-like conversation",
     "Extended thinking for deep analysis"],
    "More expensive than DeepSeek/Gemini",
    ["Opus 4.5 ($$$) - Maximum quality",
     "Sonnet 4.5 ($$) - Balanced",
     "Haiku 4.5 ($) - Fast & cheap"]
)

add_provider_slide("OpenAI (GPT)",
    "Most widely used, strong all-rounder",
    ["Largest ecosystem & integrations",
     "Best reasoning models (o3 series)",
     "Strong multimodal capabilities",
     "1M token context (GPT-4.1)"],
    "Premium pricing, especially GPT-5 and o3",
    ["GPT-5 ($$$) - Flagship",
     "o3 ($$$) - Complex reasoning",
     "GPT-4o ($$) - General use",
     "GPT-4o-mini ($) - Budget"]
)

add_provider_slide("Google (Gemini)",
    "Long context & Google integration",
    ["Longest context (up to 2M tokens)",
     "Native Google Workspace integration",
     "Good price-performance ratio",
     "Strong multimodal (text, image, video)"],
    "Less consistent on complex reasoning",
    ["Gemini 3 Pro ($$$) - Best quality",
     "Gemini Flash ($) - Fast & cheap",
     "Gemini 2.5 Pro ($$) - 2M context"]
)

add_provider_slide("DeepSeek",
    "Incredible value - 10-30x cheaper",
    ["Quality comparable to GPT-4/Sonnet",
     "Fraction of competitors' cost",
     "Open source (can self-host)",
     "Strong reasoning model (R1)"],
    "Chinese company (compliance considerations)",
    ["V3.2 ($) - General use",
     "R1 ($) - Complex reasoning",
     "Coder ($) - Programming"]
)

add_provider_slide("Open Source (Llama, Qwen, Mistral)",
    "Self-host for privacy & control",
    ["Full data privacy - your infrastructure",
     "No per-token API costs",
     "Customizable & fine-tunable",
     "No vendor lock-in"],
    "Requires GPUs & technical expertise",
    ["Llama 4 (Meta) - Most popular",
     "Qwen 3 (Alibaba) - 119 languages",
     "Mistral - European focus"]
)

add_section_slide("Use Case Recommendations")

add_usecase_slide("Customer Support / Chatbots", [
    {"need": "Cheapest option", "model": "DeepSeek V3.2 or Gemini Flash", "why": "10-30x cheaper, handles routine queries well"},
    {"need": "Best experience", "model": "Claude Sonnet 4.5", "why": "Natural conversation, excellent context understanding"},
    {"need": "Fastest responses", "model": "Claude Haiku 4.5", "why": "2-4x faster than other options"},
    {"need": "High volume (100k+)", "model": "Haiku or Gemini Flash", "why": "Keep costs manageable at scale"},
])

add_usecase_slide("Code Generation & Review", [
    {"need": "Best quality", "model": "Claude Opus 4.5", "why": "Industry-leading at writing, reviewing, and debugging code"},
    {"need": "Good quality, lower cost", "model": "DeepSeek Coder", "why": "338+ languages, excellent value"},
    {"need": "Large codebase", "model": "GPT-4.1 or Gemini 2.5 Pro", "why": "1-2M token context fits entire projects"},
    {"need": "IDE integration", "model": "Claude Haiku or Gemini Flash", "why": "Fast completions for real-time coding"},
])

add_usecase_slide("Document Analysis / RAG", [
    {"need": "Very long docs", "model": "Gemini 2.5 Pro (2M) or GPT-4.1", "why": "Largest context windows available"},
    {"need": "Best comprehension", "model": "Claude Opus 4.5", "why": "Excellent at understanding nuance"},
    {"need": "High volume", "model": "Gemini Flash or Haiku", "why": "Cost-effective for batch processing"},
    {"need": "Data privacy", "model": "Llama 4 Scout (self-hosted)", "why": "10M context, runs on your servers"},
])

add_usecase_slide("Research & Complex Analysis", [
    {"need": "Deep reasoning", "model": "OpenAI o3", "why": "Purpose-built for multi-step reasoning"},
    {"need": "Reasoning on budget", "model": "DeepSeek R1", "why": "Similar to o3 at ~1/20th the cost"},
    {"need": "Extended thinking", "model": "Claude Opus 4.5", "why": "Can 'think longer' on complex problems"},
    {"need": "Self-hosted", "model": "DeepSeek R1 (local)", "why": "Strong reasoning, runs locally"},
])

add_section_slide("Cost Comparison")

add_content_slide("Cost Tiers", [
    "$ BUDGET: DeepSeek, Gemini Flash, Claude Haiku",
    "      Best for: High volume, prototyping, cost-sensitive projects",
    "",
    "$$ MID-RANGE: Claude Sonnet, GPT-4o, Gemini Pro",
    "      Best for: Production applications, balanced needs",
    "",
    "$$$ PREMIUM: Claude Opus, GPT-5, OpenAI o3",
    "      Best for: Quality-critical apps, complex reasoning",
    "",
    "SELF-HOSTED: Llama 4, Qwen 3, DeepSeek",
    "      Best for: Data privacy, very high volume"
])

add_content_slide("Key Insight: DeepSeek Changes the Game", [
    "DeepSeek V3.2 offers quality comparable to GPT-4 and Claude Sonnet",
    "At 10-30x LOWER cost than competitors",
    "",
    "DeepSeek R1 matches OpenAI o3 reasoning at ~1/20th the price",
    "",
    "Consider DeepSeek as your default for cost-sensitive projects",
    "Move up to Claude/OpenAI only when quality is critical",
    "",
    "Caveat: Chinese company - evaluate compliance requirements"
])

add_content_slide("Resources for Deeper Research", [
    "Artificial Analysis (artificialanalysis.ai)",
    "      Real-time pricing and speed comparisons",
    "",
    "LMArena (lmarena.ai)",
    "      Try models side-by-side yourself",
    "",
    "WhatLLM (whatllm.org)",
    "      Weekly updated model comparisons",
    "",
    "Internal: See llm-selection-guide.md and index.html"
])

add_section_slide("References")

add_content_slide("Official Documentation Sources", [
    "Claude: docs.anthropic.com/en/docs/about-claude/models",
    "OpenAI: platform.openai.com/docs/models",
    "Gemini: ai.google.dev/gemini-api/docs/models",
    "DeepSeek: api-docs.deepseek.com",
    "",
    "Llama 4: github.com/meta-llama/llama-models",
    "Qwen 3: huggingface.co/collections/Qwen/qwen3",
    "Mistral: docs.mistral.ai/getting-started/models"
])

add_content_slide("Pricing & Benchmark Sources", [
    "Anthropic Pricing: anthropic.com/pricing",
    "OpenAI Pricing: openai.com/api/pricing",
    "Google Pricing: ai.google.dev/pricing",
    "DeepSeek Pricing: api-docs.deepseek.com/quick_start/pricing",
    "",
    "Artificial Analysis: artificialanalysis.ai/leaderboards/models",
    "LMArena Rankings: lmarena.ai",
    "Vellum Leaderboard: vellum.ai/llm-leaderboard"
])

add_content_slide("Summary: Quick Picks", [
    "General use, tight budget → DeepSeek V3.2",
    "Best overall quality → Claude Opus 4.5",
    "Coding → Claude Opus 4.5 or DeepSeek Coder",
    "Complex reasoning → OpenAI o3 or DeepSeek R1",
    "Long documents → Gemini 2.5 Pro (2M context)",
    "Speed priority → Claude Haiku or Gemini Flash",
    "Data privacy → Llama 4 or Qwen 3 (self-hosted)",
    "Google ecosystem → Gemini"
])

add_title_slide("Questions?", "Let's find the right LLM for your project")

# Save
prs.save(r"c:\Users\farid\OneDrive\Desktop\Code\Enerjisa\LLM Catalog\LLM_Selection_Guide_v2.pptx")
print("Presentation saved: LLM_Selection_Guide_v2.pptx")
